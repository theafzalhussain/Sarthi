"""
Document Tools — create PDF, Excel, PPT, and Word files.

WHY THESE TOOLS EXIST (a real failure):

A user asked to "make a PDF of this". The agent had no direct PDF tool —
only `python_chalao`. It imported `fpdf` and ran a script, but `fpdf`
was NOT installed:

    ModuleNotFoundError: No module named 'fpdf'
    FAIL: code exited with error (exit 1)

The agent stopped and asked the user "should I install the library?".
Two problems there:
  1. For every document type the agent had to hand-write a script —
     heavy, error-prone, and it failed whenever the library was missing.
  2. When a library was missing the agent got stuck.

Now every document type has a DEDICATED tool. And if the required
library is missing, the tool installs it ITSELF (once) and then does
the job. The user is never interrupted.

Tools:
    pdf_banao      -> PDF document (fpdf2)
    excel_banao    -> Excel spreadsheet (openpyxl)
    ppt_banao      -> PowerPoint presentation (python-pptx)
    word_banao     -> Word document (python-docx)

Each takes a simple content structure (title + sections/rows/slides)
and returns the file path so the agent can hand it to the user via
`file_kholo`.
"""

from __future__ import annotations

import asyncio
import importlib
import os
import sys
import time
from pathlib import Path

from ..devices.base import ActionResult
from .base import Tool, ToolContext

# Where documents are written — data/generated/documents
_OUTPUT_DIR = (
    Path(__file__).resolve().parent.parent.parent / "data" / "generated" / "documents"
)

# System folders where writing is forbidden (same as file_tools.py)
_BLOCKED_PATH_PARTS = (
    "system32", "/etc/", "/boot/", "/sys/", "/proc/", "/dev/",
    "program files", "windowsapps",
)


# ======================================================================
#  Library auto-install — the core fix in this file
# ======================================================================

# import module name -> pip package name (they sometimes differ)
_PACKAGE_FOR_MODULE = {
    "fpdf": "fpdf2",          # installs as fpdf2, imports as `fpdf`
    "openpyxl": "openpyxl",
    "pptx": "python-pptx",
    "docx": "python-docx",
    "PIL": "Pillow",
}

# Install lock — prevents two tools installing the same package at once
_install_lock = asyncio.Lock()


async def _ensure_library(module_name: str) -> str:
    """
    Make a Python library available — install it via pip if missing.

    Returns: "" if the library is ready, otherwise an error message.

    WHY: not every document library is pre-installed on a user's system.
    fpdf2 / python-pptx / python-docx are often missing. The agent used
    to fail in that case. Now we install ONCE and keep going — the user
    does not have to be interrupted mid-task.
    """
    # Already available?
    try:
        importlib.import_module(module_name)
        return ""
    except ImportError:
        pass

    package = _PACKAGE_FOR_MODULE.get(module_name, module_name)

    async with _install_lock:
        # Someone else may have installed it while we waited for the lock
        try:
            importlib.import_module(module_name)
            return ""
        except ImportError:
            pass

        # Run pip install (using the current interpreter's pip)
        try:
            process = await asyncio.create_subprocess_exec(
                sys.executable, "-m", "pip", "install", "--quiet", package,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            try:
                _, stderr = await asyncio.wait_for(
                    process.communicate(), timeout=180
                )
            except asyncio.TimeoutError:
                try:
                    process.kill()
                except Exception:  # noqa: BLE001
                    pass
                return (
                    f"Timed out installing '{package}' (3 min). "
                    f"Install it manually: pip install {package}"
                )
        except Exception as exc:  # noqa: BLE001
            return f"Could not install '{package}': {exc}"

        if process.returncode != 0:
            err = (stderr or b"").decode("utf-8", errors="replace").strip()
            return (
                f"Failed to install '{package}' (exit {process.returncode}).\n"
                f"{err[:400]}\n"
                f"Install it manually: pip install {package}"
            )

        # Refresh the import cache after install and try again
        importlib.invalidate_caches()
        try:
            importlib.import_module(module_name)
            return ""
        except ImportError as exc:
            return f"'{package}' installed but could not be imported: {exc}"


# ======================================================================
#  Path helpers
# ======================================================================


def _resolve_output_path(path_text: str | None, default_name: str, ext: str) -> Path:
    """
    Decide the output path.

    - if a path is given, use it (resolving env vars and ~)
    - otherwise write to data/generated/documents/ with a default name
    - append the extension if missing
    """
    if path_text and str(path_text).strip():
        expanded = os.path.expandvars(str(path_text).strip().strip('"').strip("'"))
        target = Path(expanded).expanduser()
        if not target.suffix:
            target = target.with_suffix(f".{ext}")
        return target

    _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    safe = "".join(
        c if c.isalnum() or c in " _-" else "" for c in (default_name or "document")[:50]
    ).strip().replace(" ", "_") or "document"
    return _OUTPUT_DIR / f"{safe}_{int(time.time())}.{ext}"


def _path_problem(path: Path) -> str:
    """Writing to system folders is forbidden (same check as file_tools.py)."""
    lowered = str(path).lower().replace("\\", "/")
    for blocked in _BLOCKED_PATH_PARTS:
        if blocked in lowered:
            return (
                f"'{path}' is a system location — writing there is not allowed. "
                f"Write to Desktop or Documents, or leave the path empty."
            )
    return ""


def _prepare_target(path_text, default_name, ext) -> tuple[Path | None, str]:
    """Resolve path + safety check + create parent dir. Returns (path, error)."""
    target = _resolve_output_path(path_text, default_name, ext)
    problem = _path_problem(target)
    if problem:
        return None, problem
    try:
        target.parent.mkdir(parents=True, exist_ok=True)
    except OSError as exc:
        return None, f"Could not create folder: {exc}"
    return target, ""


# ======================================================================
#  PDF
# ======================================================================


class PdfCreateTool(Tool):
    name = "pdf_banao"
    description = (
        "Create a PDF file. Use this whenever the user asks for a PDF "
        "(notes, report, resume, letter, any document). If the library "
        "is missing it installs it ITSELF — do NOT try to hack it "
        "together with fpdf via python_chalao. Provide the content as a "
        "title plus a list of sections. Each section has a heading and a "
        "body. No markdown needed — pass plain text."
    )
    parameters = {
        "type": "object",
        "properties": {
            "title": {
                "type": "string",
                "description": "Document title (shown large on the first page)",
            },
            "sections": {
                "type": "array",
                "description": (
                    "List of sections. Each section: "
                    '{"heading": "...", "body": "..."}. '
                    "Leave heading empty if you only need a body."
                ),
                "items": {
                    "type": "object",
                    "properties": {
                        "heading": {"type": "string"},
                        "body": {"type": "string"},
                    },
                },
            },
            "path": {
                "type": "string",
                "description": (
                    "Where to save it (e.g. '~/Desktop/notes.pdf'). "
                    "If left empty it is created under data/generated/documents."
                ),
            },
        },
        "required": ["title", "sections"],
    }
    risky = True

    async def run(
        self, ctx: ToolContext, title: str, sections, path: str = ""
    ) -> ActionResult:
        err = await _ensure_library("fpdf")
        if err:
            return ActionResult.failure(err)

        from fpdf import FPDF  # type: ignore

        target, problem = _prepare_target(path, title, "pdf")
        if problem:
            return ActionResult.failure(problem)

        # normalize sections — accept a list of dicts or a list of strings
        norm = _normalize_sections(sections)

        def _clean(text: str) -> str:
            # fpdf core fonts only support latin-1 — emoji/Devanagari
            # would crash. Strip unsupported characters.
            return "".join(ch if ord(ch) < 256 else "?" for ch in str(text))

        try:
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.add_page()

            pdf.set_font("Helvetica", "B", 18)
            pdf.multi_cell(0, 10, _clean(title))
            pdf.ln(4)

            for sec in norm:
                heading = sec.get("heading", "").strip()
                body = sec.get("body", "").strip()
                if heading:
                    pdf.set_font("Helvetica", "B", 13)
                    pdf.multi_cell(0, 8, _clean(heading))
                    pdf.ln(1)
                if body:
                    pdf.set_font("Helvetica", "", 11)
                    pdf.multi_cell(0, 6, _clean(body))
                    pdf.ln(3)

            pdf.output(str(target))
        except Exception as exc:  # noqa: BLE001
            return ActionResult.failure(f"Error while creating PDF: {exc}")

        size = _safe_size(target)
        return ActionResult.success(
            f"PDF created: {target}\n"
            f"  {len(norm)} sections, {size} bytes\n"
            f"  (use file_kholo to show it to the user)",
            path=str(target),
            file_path=str(target),
        )


# ======================================================================
#  EXCEL
# ======================================================================


class ExcelCreateTool(Tool):
    name = "excel_banao"
    description = (
        "Create an Excel spreadsheet (.xlsx). Use this whenever the user "
        "asks for an Excel/sheet/table (marksheet, data, budget, list). "
        "If the library is missing it installs it ITSELF. Provide the data "
        "as a list of rows — the first row is the header, the rest is data. "
        "Numbers are stored as numbers automatically."
    )
    parameters = {
        "type": "object",
        "properties": {
            "rows": {
                "type": "array",
                "description": (
                    "List of rows. Each row is a list of cells. "
                    'Example: [["Name","Marks"],["Ram",85],["Shyam",90]]. '
                    "The first row is treated as the header (bold)."
                ),
                "items": {"type": "array"},
            },
            "sheet_name": {
                "type": "string",
                "description": "Sheet name (default 'Sheet1')",
            },
            "path": {
                "type": "string",
                "description": (
                    "Where to save it. If left empty it is created under "
                    "data/generated/documents."
                ),
            },
        },
        "required": ["rows"],
    }
    risky = True

    async def run(
        self, ctx: ToolContext, rows, sheet_name: str = "Sheet1", path: str = ""
    ) -> ActionResult:
        err = await _ensure_library("openpyxl")
        if err:
            return ActionResult.failure(err)

        import openpyxl  # type: ignore
        from openpyxl.styles import Font  # type: ignore

        if not isinstance(rows, list) or not rows:
            return ActionResult.failure(
                "rows is empty — provide at least one row "
                '(e.g. [["Name","Marks"],["Ram",85]])'
            )

        target, problem = _prepare_target(path, sheet_name or "sheet", "xlsx")
        if problem:
            return ActionResult.failure(problem)

        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = (sheet_name or "Sheet1")[:31]  # Excel limit

            for r_idx, row in enumerate(rows, start=1):
                # a row may be a bare string — treat it as a single cell
                cells = row if isinstance(row, (list, tuple)) else [row]
                for c_idx, value in enumerate(cells, start=1):
                    cell = ws.cell(row=r_idx, column=c_idx, value=value)
                    if r_idx == 1:
                        cell.font = Font(bold=True)

            # Auto-fit column widths (roughly) for readability
            for column_cells in ws.columns:
                length = max(
                    (len(str(c.value)) if c.value is not None else 0)
                    for c in column_cells
                )
                col_letter = column_cells[0].column_letter
                ws.column_dimensions[col_letter].width = min(max(length + 2, 8), 50)

            wb.save(str(target))
        except Exception as exc:  # noqa: BLE001
            return ActionResult.failure(f"Error while creating Excel: {exc}")

        size = _safe_size(target)
        return ActionResult.success(
            f"Excel created: {target}\n"
            f"  {len(rows)} rows, {size} bytes\n"
            f"  (use file_kholo to show it to the user)",
            path=str(target),
            file_path=str(target),
        )


# ======================================================================
#  POWERPOINT
# ======================================================================


class PptCreateTool(Tool):
    name = "ppt_banao"
    description = (
        "Create a PowerPoint presentation (.pptx). Use this whenever the "
        "user asks for a PPT/slides/presentation. If the library is "
        "missing it installs it ITSELF. Provide a list of slides — each "
        "slide has a title and bullet points."
    )
    parameters = {
        "type": "object",
        "properties": {
            "title": {
                "type": "string",
                "description": "Presentation title (the first slide)",
            },
            "slides": {
                "type": "array",
                "description": (
                    "List of slides. Each slide: "
                    '{"title": "...", "bullets": ["point 1", "point 2"]}.'
                ),
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "bullets": {"type": "array", "items": {"type": "string"}},
                    },
                },
            },
            "path": {
                "type": "string",
                "description": (
                    "Where to save it. If left empty it is created under "
                    "data/generated/documents."
                ),
            },
        },
        "required": ["title", "slides"],
    }
    risky = True

    async def run(
        self, ctx: ToolContext, title: str, slides, path: str = ""
    ) -> ActionResult:
        err = await _ensure_library("pptx")
        if err:
            return ActionResult.failure(err)

        from pptx import Presentation  # type: ignore
        from pptx.util import Pt  # type: ignore

        target, problem = _prepare_target(path, title, "pptx")
        if problem:
            return ActionResult.failure(problem)

        if not isinstance(slides, list):
            slides = []

        try:
            prs = Presentation()

            # Title slide
            title_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = str(title)

            # Content slides
            content_layout = prs.slide_layouts[1]
            for sd in slides:
                if isinstance(sd, str):
                    sd = {"title": sd, "bullets": []}
                if not isinstance(sd, dict):
                    continue

                s = prs.slides.add_slide(content_layout)
                s.shapes.title.text = str(sd.get("title", ""))

                bullets = sd.get("bullets") or []
                if isinstance(bullets, str):
                    bullets = [bullets]

                body = s.placeholders[1].text_frame
                body.clear()
                for i, point in enumerate(bullets):
                    para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                    para.text = str(point)
                    para.font.size = Pt(18)

            prs.save(str(target))
        except Exception as exc:  # noqa: BLE001
            return ActionResult.failure(f"Error while creating PPT: {exc}")

        size = _safe_size(target)
        return ActionResult.success(
            f"Presentation created: {target}\n"
            f"  {len(slides) + 1} slides, {size} bytes\n"
            f"  (use file_kholo to show it to the user)",
            path=str(target),
            file_path=str(target),
        )


# ======================================================================
#  WORD
# ======================================================================


class WordCreateTool(Tool):
    name = "word_banao"
    description = (
        "Create a Word document (.docx). Use this whenever the user asks "
        "for a Word/doc/letter/resume/report (with formatting). If the "
        "library is missing it installs it ITSELF. Provide the content as "
        "a title plus sections (like pdf_banao)."
    )
    parameters = {
        "type": "object",
        "properties": {
            "title": {"type": "string", "description": "Document title"},
            "sections": {
                "type": "array",
                "description": (
                    'List of sections. Each section: {"heading": "...", '
                    '"body": "..."}.'
                ),
                "items": {
                    "type": "object",
                    "properties": {
                        "heading": {"type": "string"},
                        "body": {"type": "string"},
                    },
                },
            },
            "path": {
                "type": "string",
                "description": (
                    "Where to save it. If left empty it is created under "
                    "data/generated/documents."
                ),
            },
        },
        "required": ["title", "sections"],
    }
    risky = True

    async def run(
        self, ctx: ToolContext, title: str, sections, path: str = ""
    ) -> ActionResult:
        err = await _ensure_library("docx")
        if err:
            return ActionResult.failure(err)

        import docx  # type: ignore

        target, problem = _prepare_target(path, title, "docx")
        if problem:
            return ActionResult.failure(problem)

        norm = _normalize_sections(sections)

        try:
            document = docx.Document()
            document.add_heading(str(title), level=0)

            for sec in norm:
                heading = sec.get("heading", "").strip()
                body = sec.get("body", "").strip()
                if heading:
                    document.add_heading(heading, level=1)
                if body:
                    document.add_paragraph(body)

            document.save(str(target))
        except Exception as exc:  # noqa: BLE001
            return ActionResult.failure(f"Error while creating Word document: {exc}")

        size = _safe_size(target)
        return ActionResult.success(
            f"Word document created: {target}\n"
            f"  {len(norm)} sections, {size} bytes\n"
            f"  (use file_kholo to show it to the user)",
            path=str(target),
            file_path=str(target),
        )


# ======================================================================
#  Shared helpers
# ======================================================================


def _normalize_sections(sections) -> list[dict]:
    """
    Normalize sections into a list of {heading, body} dicts.

    The LLM sometimes sends a list of strings, sometimes a list of dicts,
    sometimes a single string. Handle all of them so we never crash.
    """
    if sections is None:
        return []
    if isinstance(sections, str):
        return [{"heading": "", "body": sections}]
    if not isinstance(sections, (list, tuple)):
        return [{"heading": "", "body": str(sections)}]

    result: list[dict] = []
    for item in sections:
        if isinstance(item, dict):
            result.append(
                {
                    "heading": str(item.get("heading", "") or ""),
                    "body": str(item.get("body", "") or ""),
                }
            )
        elif isinstance(item, str):
            result.append({"heading": "", "body": item})
        else:
            result.append({"heading": "", "body": str(item)})
    return result


def _safe_size(path: Path) -> int:
    try:
        return path.stat().st_size
    except OSError:
        return 0


# ======================================================================
#  Factory
# ======================================================================


def document_tools() -> list[Tool]:
    """Tools for creating PDF, Excel, PPT, and Word documents."""
    return [
        PdfCreateTool(),
        ExcelCreateTool(),
        PptCreateTool(),
        WordCreateTool(),
    ]
